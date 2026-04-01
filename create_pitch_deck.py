#!/usr/bin/env python3
"""Generate OnlyKrida Pitch Deck as PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
DARK_BG = RGBColor(0x12, 0x12, 0x12)
CARD_BG = RGBColor(0x1A, 0x1A, 0x1A)
GREEN = RGBColor(0x30, 0xD1, 0x58)
ORANGE = RGBColor(0xFF, 0x9F, 0x0A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
DIM_GRAY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color=BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=16, color=WHITE, bold=False,
             alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_accent_line(slide, left, top, width=Inches(1.5), color=GREEN):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

add_accent_line(slide, Inches(1), Inches(2.2), Inches(2), GREEN)

add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
            "ONLYKRIDA", font_size=60, color=GREEN, bold=True)

add_textbox(slide, Inches(1), Inches(3.8), Inches(10), Inches(0.8),
            "India's First AI-Powered Sports Talent Discovery Platform",
            font_size=28, color=WHITE, bold=False)

add_textbox(slide, Inches(1), Inches(4.8), Inches(10), Inches(0.6),
            "Making Every Athlete Visible. Every Scout Smarter.",
            font_size=20, color=ORANGE, bold=True)

add_textbox(slide, Inches(1), Inches(6.2), Inches(10), Inches(0.5),
            "Pitch Deck  |  2026", font_size=14, color=MID_GRAY)


# ============================================================
# SLIDE 2 — THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), ORANGE)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "THE PROBLEM", font_size=36, color=WHITE, bold=True)

problems = [
    ("500M+", "Indians play sports regularly — less than 1% ever get scouted or discovered"),
    ("Zero Platform", "No centralized system connecting athletes with scouts, coaches, and teams"),
    ("Broken Process", "Talent discovery relies on word-of-mouth, WhatsApp groups, and local tournaments"),
    ("Scout Burnout", "Scouts travel thousands of kilometers to find ONE talented player"),
    ("Invisible Athletes", "Millions of skilled athletes have no way to showcase their talent beyond their local community"),
    ("Academy Isolation", "Academies can't reach talent outside their city — limited to foot traffic and referrals"),
]

for i, (stat, desc) in enumerate(problems):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6)
    top = Inches(1.8 + row * 1.7)

    card = add_rounded_rect(slide, left, top, Inches(5.5), Inches(1.4))

    add_textbox(slide, left + Inches(0.3), top + Inches(0.15), Inches(5), Inches(0.5),
                stat, font_size=24, color=ORANGE, bold=True)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.65), Inches(4.8), Inches(0.6),
                desc, font_size=14, color=LIGHT_GRAY)


# ============================================================
# SLIDE 3 — THE SOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), GREEN)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "THE SOLUTION", font_size=36, color=WHITE, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
            "OnlyKrida is a mobile-first platform that connects the entire sports ecosystem",
            font_size=18, color=LIGHT_GRAY)

solutions = [
    ("AI-POWERED MATCHING", "Weighted scoring algorithm matches scouts to athletes based on skill, speed, stamina & position. Auto-recomputes in real-time."),
    ("VIDEO HIGHLIGHTS", "Athletes upload training clips, match highlights & achievement reels. Scouts discover talent from anywhere."),
    ("OPPORTUNITY MARKETPLACE", "Browse & apply to tryouts, tournaments, scholarships, sponsorships, and contracts — all in one place."),
    ("REAL-TIME MESSAGING", "1:1 and group chats connect athletes directly with scouts, coaches, and teams. No middlemen."),
    ("VERIFICATION PIPELINE", "School → District → State → National → Professional badges build trust and credibility."),
    ("9-ROLE ECOSYSTEM", "Athlete, Coach, Scout, Team, Academy, Trainer, Gym, Brand, Fan — every stakeholder has a dedicated experience."),
]

for i, (title, desc) in enumerate(solutions):
    row = i // 3
    col = i % 3
    left = Inches(0.5 + col * 4.1)
    top = Inches(2.4 + row * 2.4)

    card = add_rounded_rect(slide, left, top, Inches(3.8), Inches(2.0))

    add_textbox(slide, left + Inches(0.25), top + Inches(0.15), Inches(3.3), Inches(0.5),
                title, font_size=13, color=GREEN, bold=True)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.65), Inches(3.3), Inches(1.2),
                desc, font_size=12, color=LIGHT_GRAY)


# ============================================================
# SLIDE 4 — HOW ATHLETES BECOME VISIBLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), GREEN)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "HOW ATHLETES BECOME VISIBLE", font_size=36, color=WHITE, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5),
            "From invisible to discovered — the OnlyKrida athlete journey",
            font_size=18, color=LIGHT_GRAY)

steps = [
    ("01", "CREATE PROFILE", "Sign up, choose sport & position, add stats (skill, speed, stamina), set achievements, upload avatar & cover photo"),
    ("02", "UPLOAD HIGHLIGHTS", "Record training sessions, match clips, achievement moments. Video compression handles slow networks."),
    ("03", "GET AI-MATCHED", "Our algorithm scores athlete-scout fit (0-100). Scouts see top matches based on their configured preferences."),
    ("04", "RECEIVE INTEREST", "When scouts express interest, athletes get notified. See which organizations are watching — with fit scores."),
    ("05", "APPLY TO OPPORTUNITIES", "Browse tryouts, tournaments, scholarships. Apply directly. Track status: Pending → Under Review → Selected."),
    ("06", "GET DISCOVERED", "Chat directly with scouts/coaches. Get shortlisted. Move from local field to professional stage."),
]

for i, (num, title, desc) in enumerate(steps):
    row = i // 3
    col = i % 3
    left = Inches(0.5 + col * 4.1)
    top = Inches(2.3 + row * 2.5)

    card = add_rounded_rect(slide, left, top, Inches(3.8), Inches(2.1))

    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.15), Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLACK
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False

    add_textbox(slide, left + Inches(0.85), top + Inches(0.2), Inches(2.8), Inches(0.4),
                title, font_size=14, color=WHITE, bold=True)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.75), Inches(3.3), Inches(1.2),
                desc, font_size=11, color=LIGHT_GRAY)


# ============================================================
# SLIDE 5 — CORE FEATURES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), GREEN)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "CORE FEATURES — WHAT'S BUILT", font_size=36, color=WHITE, bold=True)

features = [
    ("SOCIAL FEED", "Instagram-style feed with posts (highlights, training, matches, achievements), likes, comments, follows, real-time updates"),
    ("DISCOVER", "Search athletes/scouts by sport, role, location, verification. Smart filters auto-initialize based on user role"),
    ("AI SCOUTING", "Weighted fit scoring (skill/speed/stamina/position), configurable weights per scout, real-time recomputation, shortlisting with notes"),
    ("OPPORTUNITIES", "5 types: tryouts, tournaments, sponsorships, scholarships, contracts. Full application pipeline with status tracking"),
    ("MESSAGING", "Real-time 1:1 DMs + group chats. Read receipts, unread badges, media attachments. Direct message from any profile"),
    ("TEAM DASHBOARD", "For scouts/coaches/teams: manage applications, view interested athletes, track pipeline (pending/accepted/rejected)"),
    ("ANALYTICS", "12+ tracked events: screen views, post interactions, applications, messages, follows, interest expressions"),
    ("NOTIFICATIONS", "11 types: likes, follows, comments, messages, opportunities, applications, mentions — all real-time via Supabase"),
]

for i, (title, desc) in enumerate(features):
    row = i // 4
    col = i % 4
    left = Inches(0.4 + col * 3.15)
    top = Inches(1.8 + row * 2.7)

    card = add_rounded_rect(slide, left, top, Inches(2.95), Inches(2.3))

    add_textbox(slide, left + Inches(0.15), top + Inches(0.15), Inches(2.6), Inches(0.4),
                title, font_size=12, color=GREEN, bold=True)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.6), Inches(2.6), Inches(1.5),
                desc, font_size=10, color=LIGHT_GRAY)


# ============================================================
# SLIDE 6 — AI SCOUTING DEEP DIVE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), GREEN)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "AI SCOUTING ENGINE", font_size=36, color=WHITE, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5),
            "Intelligent matching that gets smarter over time",
            font_size=18, color=LIGHT_GRAY)

# Left side - How it works
card = add_rounded_rect(slide, Inches(0.5), Inches(2.3), Inches(5.8), Inches(4.5))

txBox = add_textbox(slide, Inches(0.8), Inches(2.5), Inches(5.2), Inches(0.4),
                    "HOW IT WORKS", font_size=16, color=GREEN, bold=True)

algo_steps = [
    "1. Athletes enter performance stats (skill, speed, stamina: 0-100)",
    "2. Scouts configure matching weights (what matters most to them)",
    "3. Algorithm computes fit score: weighted combination of 4 dimensions",
    "4. Position matching: exact=100, fuzzy=80, none=20",
    "5. Top 20 matches served per scout, auto-recomputed in real-time",
    "6. Scouts shortlist, express interest — athletes get notified",
    "",
    "Formula:",
    "fit = (skill x W1 + speed x W2 + stamina x W3 + posMatch x W4) / sum(W)",
    "",
    "Default weights: Skill 35% | Speed 25% | Stamina 20% | Position 20%",
]

tf = txBox.text_frame
for step in algo_steps:
    color = ORANGE if step.startswith("fit =") or step.startswith("Default") else LIGHT_GRAY
    sz = 13 if step.startswith("fit =") else 12
    add_para(tf, step, font_size=sz, color=color, bold=step.startswith("fit ="))

# Right side - Future AI roadmap
card2 = add_rounded_rect(slide, Inches(6.8), Inches(2.3), Inches(5.8), Inches(4.5))

txBox2 = add_textbox(slide, Inches(7.1), Inches(2.5), Inches(5.2), Inches(0.4),
                     "AI ROADMAP", font_size=16, color=ORANGE, bold=True)

future = [
    "NEAR-TERM (3-6 months):",
    "  Cosine similarity for high-dimensional matching",
    "  Sport-specific stat models (not just generic 3 stats)",
    "  AI-generated scout reports via Claude API",
    "",
    "MID-TERM (6-12 months):",
    "  Video auto-tagging with computer vision",
    "  Collaborative filtering (similar scouts = similar picks)",
    "  Recency weighting (recent performance > old)",
    "",
    "LONG-TERM (12-24 months):",
    "  Embeddings (pgvector) for semantic matching",
    "  Wearable data integration (GPS, heart rate, sprint speed)",
    "  Performance prediction models",
    "  Automated highlight reel generation from match footage",
]

tf2 = txBox2.text_frame
for line in future:
    is_header = line.endswith(":")
    add_para(tf2, line, font_size=12,
             color=GREEN if is_header else LIGHT_GRAY,
             bold=is_header)


# ============================================================
# SLIDE 7 — MARKET OPPORTUNITY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), GREEN)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "MARKET OPPORTUNITY", font_size=36, color=WHITE, bold=True)

# TAM SAM SOM
markets = [
    ("TAM", "$130B+", "India Sports Market by 2030\n(KPMG/CII)", "All athletes, coaches, scouts,\nteams, academies in India"),
    ("SAM", "250M+", "Youth aged 15-30 playing\norganized/semi-organized sports", "Football, cricket, basketball,\nathletics, kabaddi in Tier 1/2 cities"),
    ("SOM", "50K+", "Active football players\nin Bengaluru (launch market)", "Starting with Football + Bengaluru\nExpanding to 8-10 sports, 10+ cities"),
]

for i, (label, num, sub1, sub2) in enumerate(markets):
    left = Inches(0.5 + i * 4.2)
    card = add_rounded_rect(slide, left, Inches(1.8), Inches(3.8), Inches(4.8))

    # Circle with label
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(1.3), Inches(2.0), Inches(1.2), Inches(1.2)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN if i == 0 else (ORANGE if i == 1 else WHITE)
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLACK
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, left + Inches(0.2), Inches(3.4), Inches(3.4), Inches(0.6),
                num, font_size=32, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), Inches(4.1), Inches(3.4), Inches(0.8),
                sub1, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), Inches(5.1), Inches(3.4), Inches(0.8),
                sub2, font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
            "No direct competitor in India combining social + AI scouting + opportunities in one mobile app",
            font_size=14, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 8 — REVENUE MODEL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), GREEN)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "REVENUE MODEL", font_size=36, color=WHITE, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
            "Athletes ALWAYS free. Revenue from organizations, premium tools, and marketplace.",
            font_size=16, color=LIGHT_GRAY)

revenue_streams = [
    ("Academy Pro", "INR 1-5K/mo", "Month 4-6", "Branded page, athlete dashboard, scout visibility"),
    ("Scout Pro (Individual)", "INR 2-5K/mo", "Month 6-9", "Advanced matching, unlimited shortlists, AI reports"),
    ("Scout Pro (Institutional)", "INR 10-25K/mo", "Month 6-9", "Team accounts, bulk scouting, analytics"),
    ("Opportunity Listings", "INR 2-5K/listing", "Month 9-12", "Post tryouts, tournaments, sponsorships"),
    ("Featured Placement", "INR 5-10K/week", "Month 9-12", "Priority visibility in athlete feeds"),
    ("Athlete Pro", "INR 199-499/mo", "Month 12-18", "Advanced analytics, performance insights"),
    ("Tournament SaaS", "INR 10-25K/event", "Month 12-18", "End-to-end tournament management"),
    ("B2B Data", "Custom pricing", "Month 12-18", "Anonymized sports talent data & insights"),
]

# Header row
header_y = Inches(2.2)
add_textbox(slide, Inches(0.5), header_y, Inches(3), Inches(0.35),
            "STREAM", font_size=11, color=GREEN, bold=True)
add_textbox(slide, Inches(3.5), header_y, Inches(2), Inches(0.35),
            "PRICE", font_size=11, color=GREEN, bold=True)
add_textbox(slide, Inches(5.5), header_y, Inches(1.8), Inches(0.35),
            "TIMELINE", font_size=11, color=GREEN, bold=True)
add_textbox(slide, Inches(7.3), header_y, Inches(5.5), Inches(0.35),
            "VALUE PROPOSITION", font_size=11, color=GREEN, bold=True)

# Divider
add_accent_line(slide, Inches(0.5), Inches(2.55), Inches(12.3), DIM_GRAY)

for i, (stream, price, timeline, value) in enumerate(revenue_streams):
    y = Inches(2.7 + i * 0.5)
    bg_color = CARD_BG if i % 2 == 0 else DARK_BG
    add_rounded_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.45), bg_color)
    add_textbox(slide, Inches(0.5), y + Inches(0.05), Inches(3), Inches(0.35),
                stream, font_size=12, color=WHITE, bold=True)
    add_textbox(slide, Inches(3.5), y + Inches(0.05), Inches(2), Inches(0.35),
                price, font_size=12, color=ORANGE, bold=True)
    add_textbox(slide, Inches(5.5), y + Inches(0.05), Inches(1.8), Inches(0.35),
                timeline, font_size=11, color=MID_GRAY)
    add_textbox(slide, Inches(7.3), y + Inches(0.05), Inches(5.5), Inches(0.35),
                value, font_size=11, color=LIGHT_GRAY)

# Target
card = add_rounded_rect(slide, Inches(3), Inches(6.8), Inches(7.3), Inches(0.5), GREEN)
add_textbox(slide, Inches(3.2), Inches(6.83), Inches(6.9), Inches(0.4),
            "MONTH 18 TARGET:  INR 5,00,000+/month  (INR 60L+/year)",
            font_size=16, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9 — GO-TO-MARKET
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), GREEN)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "GO-TO-MARKET STRATEGY", font_size=36, color=WHITE, bold=True)

phases = [
    ("PHASE 1", "Weeks 1-4", "MVP LAUNCH",
     ["Football + Bengaluru focus",
      "Seed: 100 athletes, 10 scouts, 5 academies",
      "KPIs: 50+ videos, 5+ shortlists, 10+ messages",
      "Channel: Academy partnerships + WhatsApp"]),
    ("PHASE 2", "Month 2-4", "PRODUCT-MARKET FIT",
     ["Scale to 1,000 athletes, 50 scouts",
      "Add training logs, sports passports",
      "WhatsApp deep links for viral sharing",
      "#OnlyKridaChallenge on Instagram"]),
    ("PHASE 3", "Month 4-8", "AI & ANALYTICS",
     ["10,000 athletes, 200 scouts",
      "Enhanced AI matching + video auto-tagging",
      "Wearable device integration",
      "State-level sports association partnerships"]),
    ("PHASE 4", "Month 8-18", "MULTI-SPORT SCALE",
     ["100,000 athletes across 8-10 sports",
      "10+ cities + Dubai expansion",
      "7 Indian languages (vernacular support)",
      "National federation & IPL/ISL partnerships"]),
]

for i, (phase, time, title, bullets) in enumerate(phases):
    left = Inches(0.3 + i * 3.2)
    card = add_rounded_rect(slide, left, Inches(1.8), Inches(3.0), Inches(5.2))

    # Phase badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.15), Inches(1.95), Inches(1.2), Inches(0.35)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = GREEN
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = phase
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLACK
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, left + Inches(1.5), Inches(1.95), Inches(1.3), Inches(0.35),
                time, font_size=10, color=MID_GRAY)

    add_textbox(slide, left + Inches(0.2), Inches(2.5), Inches(2.6), Inches(0.4),
                title, font_size=15, color=WHITE, bold=True)

    txBox = add_textbox(slide, left + Inches(0.2), Inches(3.1), Inches(2.6), Inches(3.5),
                        "", font_size=12, color=LIGHT_GRAY)
    tf = txBox.text_frame
    for j, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Calibri"
        p.space_before = Pt(8)


# ============================================================
# SLIDE 10 — WHAT'S BUILT (STATUS)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), GREEN)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "WHAT'S ALREADY BUILT", font_size=36, color=WHITE, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
            "Production-ready. TypeScript strict. 0 compile errors. Supabase backend with RLS security.",
            font_size=16, color=LIGHT_GRAY)

built_items = [
    "9-role authentication with dedicated signup flows",
    "AI scouting engine with weighted fit scoring (0-100)",
    "Real-time 1:1 + group messaging with read receipts",
    "Opportunity marketplace with full application pipeline",
    "Social feed: posts, likes, comments, follows (real-time)",
    "Discover with sport/role/location/verified filters",
    "Video & image upload with compression",
    "Notifications center (11 types, all real-time)",
    "Team dashboard with application management",
    "Scout preferences + shortlisting with notes",
    "Analytics tracking (12+ events, batched flush)",
    "Profile customization (cover photo, avatar, stats, achievements)",
    "Verification badge system (5-tier pipeline)",
    "16+ database tables with triggers & indexes",
    "Row Level Security on all tables",
    "Offline caching via AsyncStorage",
]

for i, item in enumerate(built_items):
    row = i // 2
    col = i % 2
    left = Inches(0.5 + col * 6.3)
    y = Inches(2.1 + row * 0.6)

    # Green check circle
    check = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, y + Inches(0.05), Inches(0.25), Inches(0.25)
    )
    check.fill.solid()
    check.fill.fore_color.rgb = GREEN
    check.line.fill.background()

    add_textbox(slide, left + Inches(0.35), y, Inches(5.5), Inches(0.35),
                item, font_size=12, color=WHITE)


# ============================================================
# SLIDE 11 — TECH STACK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), GREEN)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "TECHNOLOGY STACK", font_size=36, color=WHITE, bold=True)

stack_groups = [
    ("FRONTEND", [
        "React Native 0.81 + Expo SDK 54",
        "TypeScript 5.9 (strict mode)",
        "Expo Router (file-based routing)",
        "Lucide Icons + Expo Blur/Gradient",
    ]),
    ("BACKEND", [
        "Supabase (PostgreSQL + Auth)",
        "Supabase Realtime (subscriptions)",
        "Supabase Storage (media uploads)",
        "Row Level Security on all tables",
    ]),
    ("STATE", [
        "10 Context Providers",
        "AsyncStorage (offline caching)",
        "React Query (data fetching)",
        "Zustand (available, secondary)",
    ]),
    ("AI / DATA", [
        "Custom weighted scoring algorithm",
        "Real-time recomputation engine",
        "Analytics event pipeline",
        "Ready for ML/embeddings upgrade",
    ]),
]

for i, (group, items) in enumerate(stack_groups):
    left = Inches(0.3 + i * 3.2)
    card = add_rounded_rect(slide, left, Inches(1.8), Inches(3.0), Inches(4.0))

    add_textbox(slide, left + Inches(0.2), Inches(2.0), Inches(2.6), Inches(0.4),
                group, font_size=14, color=GREEN, bold=True)

    txBox = add_textbox(slide, left + Inches(0.2), Inches(2.6), Inches(2.6), Inches(3.0),
                        "", font_size=12, color=LIGHT_GRAY)
    tf = txBox.text_frame
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Calibri"
        p.space_before = Pt(12)


# ============================================================
# SLIDE 12 — COMPETITIVE EDGE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), GREEN)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "WHY ONLYKRIDA WINS", font_size=36, color=WHITE, bold=True)

edges = [
    ("FIRST MOVER", "No platform in India combines social networking + AI scouting + opportunity marketplace for sports"),
    ("NETWORK EFFECTS", "More athletes attract more scouts, which attract more opportunities, which attract more athletes"),
    ("TRUST PIPELINE", "5-tier verification (School to Professional) builds credibility that no competitor offers"),
    ("MOBILE-NATIVE", "Built for India's smartphone-first, 4G reality with video compression and offline caching"),
    ("9-ROLE ECOSYSTEM", "Not just athletes — coaches, scouts, teams, academies, gyms, brands, trainers, fans all have dedicated experiences"),
    ("AI-UPGRADEABLE", "Simple weighted scoring today, ready for cosine similarity, embeddings, collaborative filtering tomorrow"),
    ("FREE FOR ATHLETES", "Core principle: never paywall discovery, search, or applying. Revenue from organizations."),
    ("LOCALIZATION READY", "Architecture supports 7 Indian languages. Built for Bengaluru, designed for all of India + Dubai"),
]

for i, (title, desc) in enumerate(edges):
    row = i // 2
    col = i % 2
    left = Inches(0.5 + col * 6.3)
    top = Inches(1.7 + row * 1.35)

    card = add_rounded_rect(slide, left, top, Inches(5.9), Inches(1.1))

    add_textbox(slide, left + Inches(0.2), top + Inches(0.1), Inches(2.5), Inches(0.35),
                title, font_size=13, color=GREEN, bold=True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.5), Inches(5.4), Inches(0.5),
                desc, font_size=11, color=LIGHT_GRAY)


# ============================================================
# SLIDE 13 — CLOSING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(4.5), Inches(1.8), Inches(4.3), GREEN)

add_textbox(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
            "ONLYKRIDA", font_size=60, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
            "Making Every Athlete Visible.", font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.8),
            "Every Scout Smarter.", font_size=28, color=ORANGE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.6),
            "500 million athletes deserve to be seen.", font_size=18, color=MID_GRAY,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
            "Let's talk.  |  onlysports.app", font_size=16, color=LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER)


# Save
output_path = "/Users/anirudhtumuluru/onlysports-platform/OnlyKrida_Pitch_Deck.pptx"
prs.save(output_path)
print(f"Pitch deck saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
