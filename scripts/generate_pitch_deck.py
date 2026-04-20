#!/usr/bin/env python3
"""
OnlyKrida Complete Pitch Deck Generator
Generates a 20-slide investor + engineering grade PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme Colors ──────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD    = RGBColor(0x1A, 0x1A, 0x2E)
GREEN      = RGBColor(0x30, 0xD1, 0x58)
ORANGE     = RGBColor(0xFF, 0x9F, 0x0A)
CYAN       = RGBColor(0x64, 0xD2, 0xFF)
RED        = RGBColor(0xFF, 0x45, 0x3A)
WHITE      = RGBColor(0xF0, 0xF0, 0xF0)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
MUTED      = RGBColor(0x88, 0x88, 0x88)
PURPLE     = RGBColor(0xBF, 0x5A, 0xF2)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txbox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, bullet_color=GREEN):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        # Bullet character
        run_b = p.add_run()
        run_b.text = "  >  "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bullet_color
        run_b.font.bold = True
        # Text
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
    return txbox

def add_table(slide, left, top, width, height, data, col_widths=None):
    """data = list of rows, first row is header."""
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w

    for r_idx, row in enumerate(data):
        for c_idx, cell_text in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            # Style
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = GREEN
                else:
                    paragraph.font.color.rgb = WHITE
                paragraph.alignment = PP_ALIGN.LEFT
            # Cell fill
            cell_fill = cell.fill
            cell_fill.solid()
            if r_idx == 0:
                cell_fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x25)
            elif r_idx % 2 == 0:
                cell_fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
            else:
                cell_fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x22)
    return table_shape

def new_slide(title_text, subtitle_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide)
    # Green accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()
    # Title
    add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.7), title_text, font_size=32, color=WHITE, bold=True)
    if subtitle_text:
        add_text(slide, Inches(0.6), Inches(0.95), Inches(10), Inches(0.5), subtitle_text, font_size=16, color=MUTED)
    # Bottom line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.1), Inches(12), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x2A)
    line.line.fill.background()
    # Footer
    add_text(slide, Inches(0.6), Inches(7.15), Inches(4), Inches(0.3), "OnlyKrida  |  Confidential", font_size=10, color=MUTED)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Large green circle accent
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(7), Inches(7))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x15, 0x30, 0x15)
circle.line.fill.background()

add_text(slide, Inches(1), Inches(1.5), Inches(8), Inches(1), "OnlyKrida", font_size=60, color=GREEN, bold=True)
add_text(slide, Inches(1), Inches(2.5), Inches(10), Inches(0.8), "India's First AI-Powered Sports Talent Discovery Platform", font_size=28, color=WHITE, bold=True)
add_text(slide, Inches(1), Inches(3.5), Inches(10), Inches(0.6), "Making Every Athlete Discoverable  |  Powered by Claude AI", font_size=18, color=LIGHT_GRAY)

# Stats bar
for i, (num, label) in enumerate([("400M+", "Sports Participants"), ("$7B", "Indian Sports Market"), ("< 1%", "Professionally Scouted"), ("9", "User Roles")]):
    x = Inches(1 + i * 2.8)
    add_shape(slide, x, Inches(4.8), Inches(2.4), Inches(1.2), BG_CARD, GREEN)
    add_text(slide, x, Inches(4.9), Inches(2.4), Inches(0.6), num, font_size=28, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(5.4), Inches(2.4), Inches(0.5), label, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(6.5), Inches(6), Inches(0.4), "Investor + Engineering Deck  |  April 2026  |  Confidential", font_size=12, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("The Problem", "Why India's sports talent ecosystem is broken")

problems = [
    ("No Digital Identity", "Athletes have no centralized profile. Talent dies in WhatsApp groups and local newspapers.", RED),
    ("Geographic Bias", "Scouts only visit metro cities. 70% of India's athletes are in Tier 2/3 cities and rural areas.", ORANGE),
    ("Pay-to-Play Culture", "Access to trials costs money. Talented athletes from low-income families are systematically excluded.", RED),
    ("Fragmented Data", "Fitness data lives on paper. No standardized, verified performance metrics exist.", ORANGE),
    ("Scout Inefficiency", "Scouts travel thousands of miles to see 10 athletes. 90% of trips yield zero signings.", RED),
    ("No Verified Credentials", "Fake achievements are rampant. There's no way to verify an athlete's actual performance.", ORANGE),
]

for i, (title, desc, color) in enumerate(problems):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.6 + row * 2.6)
    card = add_shape(slide, x, y, Inches(3.8), Inches(2.2), BG_CARD, color)
    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.5), title, font_size=18, color=color, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(1.2), desc, font_size=13, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("The Solution", "OnlyKrida — LinkedIn meets Instagram for Sports")

add_text(slide, Inches(0.6), Inches(1.6), Inches(11), Inches(0.6),
         "A hybrid platform combining professional credibility, content engagement, and AI-powered scouting",
         font_size=18, color=LIGHT_GRAY)

pillars = [
    ("Digital Sports Identity", "Verified profiles with achievements,\nfitness data, and highlight reels", GREEN),
    ("AI-Powered Discovery", "Claude AI matches athletes to scouts\nusing weighted preference algorithms", CYAN),
    ("Content Platform", "Instagram-like feed for training\nhighlights, match clips, and achievements", ORANGE),
    ("Opportunity Marketplace", "Tryouts, tournaments, scholarships,\nand contracts in one place", PURPLE),
]

for i, (title, desc, color) in enumerate(pillars):
    x = Inches(0.6 + i * 3.1)
    card = add_shape(slide, x, Inches(2.5), Inches(2.8), Inches(3.2), BG_CARD, color)
    # Icon circle
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.05), Inches(2.8), Inches(0.7), Inches(0.7))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text(slide, x + Inches(0.2), Inches(3.7), Inches(2.4), Inches(0.5), title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(4.3), Inches(2.4), Inches(1.2), desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0.6), Inches(6.2), Inches(10), Inches(0.5),
         "Free forever for athletes  |  Revenue from scout subscriptions, team dashboards, and brand sponsorships",
         font_size=14, color=GREEN, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MARKET OPPORTUNITY
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("Market Opportunity", "India's sports ecosystem is massive, underserved, and growing fast")

data = [
    ["Metric", "Value", "Source"],
    ["Total Sports Participants (India)", "400M+", "FICCI Sports Report 2025"],
    ["Indian Sports Industry Size", "$7 Billion", "KPMG India Sports Report"],
    ["Annual Growth Rate", "15% YoY", "Deloitte Sports Industry Analysis"],
    ["Professional Athletes (% of participants)", "< 1%", "Sports Authority of India"],
    ["Grassroots Athletes Without Digital Presence", "99%+", "OnlyKrida Research"],
    ["Coaches & Trainers in India", "2M+", "National Sports Census"],
    ["Active Sports Scouts (organized)", "~5,000", "ISL / IPL / Federation Data"],
    ["Sports Academies", "15,000+", "FICCI Education Report"],
    ["Annual Sports Scholarships (unfilled)", "30%+ go unfilled", "University Grants Commission"],
]
add_table(slide, Inches(0.6), Inches(1.6), Inches(11.5), Inches(4.5), data)

add_shape(slide, Inches(0.6), Inches(6.2), Inches(5.5), Inches(0.7), BG_CARD, GREEN)
add_text(slide, Inches(0.8), Inches(6.25), Inches(5), Inches(0.6),
         "TAM: $7B  |  SAM: $800M  |  SOM (Year 3): $12M", font_size=16, color=GREEN, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("How It Works", "From signup to discovery in minutes")

steps = [
    ("1", "Sign Up &\nChoose Role", "Athletes, scouts, coaches,\nteams select their role.\n5 role-specific signup forms."),
    ("2", "Build Your\nProfile", "Add sport, position, stats,\nachievements, and upload\nhighlight videos."),
    ("3", "Take Fitness\nTests", "Beep test, sprints, agility,\nvertical jump. Get your\nzone: Starter to Unstoppable."),
    ("4", "Create\nContent", "Post training clips, match\nhighlights, and achievements\nto your feed."),
    ("5", "Get\nDiscovered", "AI matches you with scouts\nbased on fit scores.\nScouts shortlist and message."),
    ("6", "Apply to\nOpportunities", "Browse tryouts, tournaments,\nscholarships, contracts.\nApply with one tap."),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.4 + i * 2.1)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), Inches(1.7), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    add_text(slide, x + Inches(0.65), Inches(1.73), Inches(0.6), Inches(0.55), num, font_size=22, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Arrow
    if i < 5:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(1.4), Inches(1.95), Inches(0.7), Pt(2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x2A)
        arrow.line.fill.background()
    # Card
    card = add_shape(slide, x, Inches(2.5), Inches(1.9), Inches(3.8), BG_CARD)
    add_text(slide, x + Inches(0.15), Inches(2.7), Inches(1.6), Inches(0.8), title, font_size=14, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(3.5), Inches(1.7), Inches(2.5), desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PRODUCT FEATURES
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("Product Features", "Everything built and working in MVP")

features = [
    ("Social Feed", "Image & video posts, likes,\ncomments, shares, infinite scroll.\n4 post types: highlight, training,\nmatch, achievement.", GREEN),
    ("AI Scouting", "Claude Opus-powered matching.\nFit scores 0-100 with breakdown.\nSmart recommendations.\nProfile coaching suggestions.", CYAN),
    ("Messaging", "Real-time 1-on-1 and group DMs.\nMedia sharing, read receipts,\npost sharing via message.\nUnread count badges.", ORANGE),
    ("Fitness Testing", "5 test types with live audio cues.\n6 growth zones (never demotivating).\nCoach verification system.\nProgress tracking & history.", PURPLE),
    ("Opportunities", "Tryouts, tournaments, scholarships,\nsponsorship, contracts.\nApply with cover letter.\nTeam-side application management.", RED),
    ("Discovery Engine", "Filter by sport, role, location,\nverified status. Trending athletes.\nTrigram text search.\n5 concurrent filter lists.", RGBColor(0xFF, 0xD6, 0x0A)),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.5 + row * 2.8)
    card = add_shape(slide, x, y, Inches(3.8), Inches(2.5), BG_CARD, color)
    add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.45), title, font_size=18, color=color, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.4), Inches(1.6), desc, font_size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AI-POWERED SCOUTING
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("AI-Powered Scouting Engine", "Claude Opus 4.6 with adaptive thinking for deep athlete analysis")

add_bullet_list(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(5), [
    "Scouts set preferences: sport, positions, regions, and weight factors (skill, speed, stamina, position match)",
    "AI analyzes all athlete profiles against scout preferences using Claude Opus with adaptive thinking",
    "Each athlete receives a fit score (0-100) with detailed breakdown by category",
    "Results cached in ai_recommendations table, refreshed every 6 hours",
    "Profile coaching: Claude Sonnet gives athletes 3-5 actionable improvement tips",
    "Opportunity matcher: AI ranks opportunities by athlete fit",
    "Athlete potential analysis: Deep trajectory prediction for long-term scouting",
    "Krida AI Assistant: Full conversational chat for questions about sports, training, careers",
], font_size=14)

# AI function table
ai_data = [
    ["AI Function", "Model", "Purpose", "Avg Latency"],
    ["Smart Scout Recommendations", "Opus + Thinking", "Rank athletes by fit score", "3-5s"],
    ["Athlete Potential Analysis", "Opus + Thinking", "Deep trajectory prediction", "4-6s"],
    ["Opportunity Matcher", "Opus", "Match athletes to opportunities", "2-4s"],
    ["Profile Coaching", "Sonnet", "Improvement suggestions", "1-2s"],
    ["Profile Summary", "Sonnet", "Scout-readable bio generation", "1-2s"],
    ["Chat Response", "Sonnet", "Conversational AI assistant", "1-2s"],
    ["Knowledge Search", "Sonnet", "RAG-style sports knowledge", "1-2s"],
]
add_table(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(3.8), ai_data)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FITNESS TESTING
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("Fitness Testing System", "Standardized, verified athletic performance data")

# Test types
test_data = [
    ["Test Type", "What It Measures", "Zones"],
    ["Yo-Yo IR1 (Beep Test)", "Cardiovascular endurance, VO2max", "6 zones: Starter to Unstoppable"],
    ["20m Sprint", "Acceleration, explosive speed", "6 zones based on time"],
    ["40m Sprint", "Top-end speed, speed endurance", "6 zones based on time"],
    ["Agility T-Test", "Change of direction, agility", "6 zones based on time"],
    ["Vertical Jump", "Lower body power, explosiveness", "6 zones based on height"],
]
add_table(slide, Inches(0.6), Inches(1.5), Inches(7), Inches(2.8), test_data)

# Verification tiers
ver_data = [
    ["Verification Tier", "How", "Trust Level"],
    ["Self-Reported", "Athlete enters manually", "Low (badge: none)"],
    ["App-Measured", "Live test with phone sensors", "Medium (badge: app)"],
    ["Coach-Verified", "Coach attests results", "High (badge: coach)"],
    ["Center-Tested", "Official testing center", "Highest (badge: verified)"],
]
add_table(slide, Inches(0.6), Inches(4.6), Inches(7), Inches(2.2), ver_data)

# Zone philosophy
add_shape(slide, Inches(8), Inches(1.5), Inches(4.8), Inches(5.3), BG_CARD, GREEN)
add_text(slide, Inches(8.2), Inches(1.7), Inches(4.4), Inches(0.5), "Zone Philosophy", font_size=18, color=GREEN, bold=True)
add_bullet_list(slide, Inches(8.2), Inches(2.3), Inches(4.4), Inches(4), [
    "Never demotivating labels",
    "No 'fail', 'poor', or 'beginner'",
    "Growth-oriented language only",
    "Starter = just beginning the journey",
    "Building = actively improving",
    "Rising = showing real potential",
    "Strong = competitive level",
    "Elite = top-tier performer",
    "Unstoppable = world-class",
], font_size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ROLE-SPECIFIC UX
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("9 User Roles, 9 Experiences", "Every role gets a tailored interface, not a generic dashboard")

roles = [
    ("Athlete", "Profile, highlights, fitness tests,\nopportunities, AI coaching", GREEN),
    ("Scout", "AI recommendations, shortlists,\npreferences, athlete comparison", ORANGE),
    ("Coach", "Team roster, athlete management,\nverify fitness tests, post content", CYAN),
    ("Team", "Create opportunities, manage\napplications, team dashboard", RED),
    ("Fan", "Follow athletes, trending feed,\nsports news, engagement", PURPLE),
    ("Trainer", "Training content, certifications,\nathlete connections", GREEN),
    ("Gym", "Facility showcase, programs,\nmembership, athlete network", ORANGE),
    ("Brand", "Sponsorship opportunities,\nathlete search, campaign tools", CYAN),
    ("Academy", "Student athletes, programs,\nrecruitment, performance tracking", RED),
]

for i, (role, desc, color) in enumerate(roles):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.5 + row * 1.85)
    card = add_shape(slide, x, y, Inches(3.8), Inches(1.6), BG_CARD, color)
    add_text(slide, x + Inches(0.2), y + Inches(0.12), Inches(3.4), Inches(0.4), role, font_size=16, color=color, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.55), Inches(3.4), Inches(0.9), desc, font_size=11, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("Technology Stack", "Modern, scalable, cost-efficient architecture")

stack_data = [
    ["Layer", "Technology", "Why"],
    ["Frontend", "React Native (Expo SDK 54)", "Cross-platform iOS/Android/Web from one codebase"],
    ["Language", "TypeScript 5.9", "Type safety across 6300+ LOC"],
    ["Routing", "Expo Router 6 (file-based)", "Next.js-style routing for React Native"],
    ["Backend", "Supabase (PostgreSQL + Auth + Storage + Realtime)", "Firebase alternative with SQL power, RLS security"],
    ["AI Engine", "Claude API (Opus 4.6 + Sonnet 4.6)", "Best-in-class reasoning for scouting analysis"],
    ["State Mgmt", "12 React Context providers + React Query", "Server state caching + local state management"],
    ["Media", "expo-image (cached) + expo-video", "Disk/memory caching, blurhash placeholders"],
    ["Icons", "Lucide React Native", "700+ consistent SVG icons"],
    ["Notifications", "expo-notifications", "Push notification infrastructure"],
    ["Build", "EAS Build + EAS Submit", "Cloud builds for iOS/Android distribution"],
    ["Search", "PostgreSQL pg_trgm", "Trigram fuzzy text search on names/locations"],
]
add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5.3), stack_data)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DATABASE ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("Database Architecture", "13 PostgreSQL tables with Row-Level Security on every table")

db_data = [
    ["Table", "Rows (est.)", "Key Columns", "Security"],
    ["profiles", "Users", "role, sport, position, achievements (JSONB), stats (JSONB), verified", "Public read, own update"],
    ["posts", "Content", "type, image_url, video_url, likes/comments/views counts", "Public read, own CRUD"],
    ["opportunities", "Jobs", "category, sport, location, deadline, requirements, compensation", "Public read, team CRUD"],
    ["applications", "Apps", "opportunity_id, athlete_id, status, cover_letter", "Applicant + owner only"],
    ["messages", "DMs", "sender_id, receiver_id, content, media_url, status, read", "Sender/receiver only"],
    ["notifications", "Alerts", "type (11 types), title, message, read, data (JSONB)", "Own user only"],
    ["follows", "Graph", "follower_id, following_id (unique, no self-follows)", "Public read, own CRUD"],
    ["likes", "Engage", "user_id, post_id (unique pair)", "Public read, own CRUD"],
    ["comments", "Engage", "user_id, post_id, content, likes_count", "Public read, own CRUD"],
    ["player_stats", "Stats", "sport, position, skill/speed/stamina (0-100)", "Public read, own CRUD"],
    ["scout_preferences", "Config", "sport, positions[], weight_skill/speed/stamina/position", "Own user CRUD"],
    ["ai_recommendations", "AI", "scout_id, player_id, fit_score, breakdown (JSONB)", "Public read, own CRUD"],
    ["comment_likes", "Engage", "user_id, comment_id (unique pair)", "Public read, own CRUD"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(12.5), Inches(5.3), db_data)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("Business Model", "Free for athletes forever — monetize the demand side")

model_data = [
    ["Revenue Stream", "Target", "Pricing", "Est. Year 1 Revenue"],
    ["Scout Basic", "Independent scouts", "Rs 999/mo ($12)", "Rs 6L ($7.2K)"],
    ["Scout Pro", "Agency scouts", "Rs 2,499/mo ($30)", "Rs 15L ($18K)"],
    ["Scout Enterprise", "Federations, ISL/IPL", "Rs 9,999/mo ($120)", "Rs 12L ($14.4K)"],
    ["Team Starter", "Local clubs", "Rs 1,999/mo ($24)", "Rs 12L ($14.4K)"],
    ["Team Pro", "Professional teams", "Rs 4,999/mo ($60)", "Rs 18L ($21.6K)"],
    ["Brand Sponsorship", "Sports brands", "Rs 25K-2L/campaign", "Rs 30L ($36K)"],
    ["Featured Opportunities", "Teams posting tryouts", "Rs 500/listing boost", "Rs 6L ($7.2K)"],
    ["Data Insights API", "Analysts, media", "Rs 4,999/mo ($60)", "Rs 6L ($7.2K) (Phase 2)"],
    ["", "", "TOTAL YEAR 1 (est.)", "Rs 1.05 Cr ($126K)"],
]
add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(4.5), model_data)

add_shape(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.7), BG_CARD, GREEN)
add_text(slide, Inches(0.8), Inches(6.25), Inches(11), Inches(0.6),
         "Key principle: Athletes never pay. The platform's value grows with athlete volume, which drives scout/team willingness to pay.",
         font_size=14, color=GREEN, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — REVENUE PROJECTIONS
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("Revenue Projections", "Conservative estimates based on Indian sports market benchmarks")

rev_data = [
    ["Metric", "Year 1", "Year 2", "Year 3"],
    ["Total Users", "50,000", "250,000", "1,000,000"],
    ["Athletes", "40,000 (80%)", "200,000", "800,000"],
    ["Scouts (paid)", "50", "300", "1,200"],
    ["Teams (paid)", "50", "200", "800"],
    ["Brands (active)", "5", "25", "100"],
    ["Monthly Recurring Revenue", "Rs 8.75L", "Rs 52.5L", "Rs 2.1 Cr"],
    ["Annual Revenue", "Rs 1.05 Cr ($126K)", "Rs 6.3 Cr ($756K)", "Rs 25.2 Cr ($3M)"],
    ["Gross Margin", "85%", "82%", "80%"],
    ["CAC (Athletes)", "Rs 50 ($0.60)", "Rs 30 ($0.36)", "Rs 15 ($0.18)"],
    ["CAC (Scouts)", "Rs 5,000 ($60)", "Rs 3,000 ($36)", "Rs 2,000 ($24)"],
    ["LTV:CAC (Scouts)", "6x", "10x", "15x"],
    ["Monthly Burn Rate", "Rs 5L ($6K)", "Rs 25L ($30K)", "Rs 80L ($96K)"],
    ["Break-even Month", "—", "Month 18", "—"],
]
add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5.3), rev_data)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — GO-TO-MARKET
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("Go-to-Market Strategy", "India first, then Dubai and global expansion")

add_bullet_list(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(2.5), [
    "Phase 1 (Months 1-6): 5 Indian cities — Mumbai, Delhi, Bangalore, Chennai, Kolkata",
    "Phase 2 (Months 7-12): Tier 2 cities — Pune, Hyderabad, Ahmedabad, Jaipur, Lucknow",
    "Phase 3 (Year 2): Dubai / UAE — Indian expat athlete community + Middle East sports market",
    "Phase 4 (Year 3): Southeast Asia — similar grassroots talent discovery gaps",
], font_size=14)

add_text(slide, Inches(0.6), Inches(4.2), Inches(5), Inches(0.5), "Growth Channels", font_size=20, color=GREEN, bold=True)
add_bullet_list(slide, Inches(0.6), Inches(4.7), Inches(5.5), Inches(2), [
    "College sports departments — bulk onboarding of entire teams",
    "Sports academies — partnership for fitness test verification",
    "Coaches as ambassadors — each coach brings 20-50 athletes",
    "Social media — athlete highlight reels go viral (TikTok/Reels/Shorts)",
    "Sports federation partnerships — official talent identification tool",
], font_size=13)

# Right side - market sizing
add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), BG_CARD, GREEN)
add_text(slide, Inches(7), Inches(1.7), Inches(5.4), Inches(0.5), "India Sports Market Snapshot", font_size=18, color=GREEN, bold=True)
india_stats = [
    "65% of India's population is under 35",
    "Cricket alone has 200M+ active players",
    "Football growing at 30% YoY in participation",
    "Kabaddi Pro League created $100M market in 5 years",
    "Government 'Khelo India' program: Rs 1,756 Cr budget",
    "IPL proved Indian sports = massive commercial opportunity",
    "Zero tech platforms serve grassroots athletes",
    "OnlyKrida is first-mover in this space",
]
add_bullet_list(slide, Inches(7), Inches(2.3), Inches(5.4), Inches(4), india_stats, font_size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — COMPETITIVE LANDSCAPE
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("Competitive Landscape", "No direct competitor serves India's grassroots athletes")

comp_data = [
    ["Feature", "OnlyKrida", "Hudl", "NCSA", "TransferMarkt", "SportsRecruits"],
    ["India Focus", "YES (primary)", "No", "No (US only)", "Partial", "No (US only)"],
    ["Free for Athletes", "YES (forever)", "Freemium", "No ($)", "Free (view)", "No ($)"],
    ["AI Scouting", "YES (Claude AI)", "Basic analytics", "No", "No", "No"],
    ["Fitness Testing", "YES (5 types)", "No", "No", "No", "No"],
    ["Social Feed", "YES (Instagram-like)", "Limited", "No", "Forum", "No"],
    ["Messaging", "YES (real-time)", "No", "Yes", "No", "Yes"],
    ["9 User Roles", "YES", "2 roles", "3 roles", "View only", "3 roles"],
    ["Opportunities", "YES (5 categories)", "No", "Limited", "Transfer only", "Scholarships"],
    ["Multi-sport", "YES (all sports)", "US football", "US college", "Football (soccer)", "US college"],
    ["Mobile App", "YES (iOS/Android/Web)", "Yes", "Yes", "Yes", "Limited"],
    ["Verification System", "YES (4 tiers)", "No", "No", "Editorial", "No"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(12.5), Inches(5.3), comp_data)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — TRACTION / MVP STATUS
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("MVP Traction & Status", "Fully functional platform — ready for user acquisition")

metrics = [
    ("47", "Components Built"),
    ("13", "Database Tables"),
    ("12", "Context Providers"),
    ("6,300+", "Lines of Code"),
    ("9", "User Roles"),
    ("7", "AI Functions"),
    ("5", "Fitness Test Types"),
    ("20+", "Screens"),
]

for i, (num, label) in enumerate(metrics):
    col = i % 4
    row = i // 4
    x = Inches(0.6 + col * 3.1)
    y = Inches(1.6 + row * 1.5)
    card = add_shape(slide, x, y, Inches(2.8), Inches(1.2), BG_CARD, GREEN)
    add_text(slide, x, y + Inches(0.1), Inches(2.8), Inches(0.6), num, font_size=28, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.65), Inches(2.8), Inches(0.4), label, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(0.5), "Features Completed in MVP", font_size=20, color=GREEN, bold=True)
add_bullet_list(slide, Inches(0.6), Inches(5.3), Inches(5.5), Inches(1.5), [
    "Complete authentication with 5 role-specific signup forms",
    "6 role-specific home dashboards",
    "Real-time messaging (1-on-1 + group)",
    "AI scouting with Claude Opus fit scores",
], font_size=13)
add_bullet_list(slide, Inches(6.5), Inches(5.3), Inches(5.5), Inches(1.5), [
    "Fitness testing with 5 types and verification",
    "Opportunity marketplace with applications",
    "Content feed with image/video support",
    "Supabase backend with RLS security on every table",
], font_size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("Product Roadmap", "3-phase plan from MVP to scale")

phases = [
    ("Phase 1: MVP (DONE)", GREEN, [
        "9-role authentication system",
        "Social feed with media uploads",
        "AI scouting (Claude Opus/Sonnet)",
        "Fitness testing (5 types, 6 zones)",
        "Real-time messaging & notifications",
        "Opportunity marketplace",
        "Discovery engine with filters",
    ]),
    ("Phase 2: Growth (Q2-Q3 2026)", ORANGE, [
        "Payment integration (Razorpay)",
        "Admin dashboard & moderation",
        "Push notifications",
        "Analytics & event tracking",
        "Performance optimization sprint",
        "Video streaming (HLS/Mux)",
        "CI/CD pipeline",
    ]),
    ("Phase 3: Scale (Q4 2026+)", CYAN, [
        "ML-powered recommendation engine",
        "Video analysis (pose detection)",
        "Live match streaming",
        "Advanced analytics dashboard",
        "International expansion (Dubai)",
        "i18n (Hindi, Tamil, Arabic)",
        "Offline-first architecture",
    ]),
]

for i, (title, color, items) in enumerate(phases):
    x = Inches(0.6 + i * 4.1)
    card = add_shape(slide, x, Inches(1.5), Inches(3.8), Inches(5.3), BG_CARD, color)
    add_text(slide, x + Inches(0.2), Inches(1.65), Inches(3.4), Inches(0.5), title, font_size=16, color=color, bold=True)
    add_bullet_list(slide, x + Inches(0.15), Inches(2.2), Inches(3.5), Inches(4.3), items, font_size=12, color=LIGHT_GRAY, bullet_color=color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — CRITICAL IMPROVEMENTS
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("Top 5 Immediate Priorities", "Critical improvements before scaling")

priorities = [
    ("P0: Secure AI API Key", "Move Anthropic API key from client to Supabase Edge Function.\nCurrently exposed in JS bundle — critical security vulnerability.", RED, "3 days"),
    ("P0: Automated Testing", "Zero tests today. Need Jest + RTL for contexts, Maestro for E2E.\nTarget 70% coverage on business logic before scaling.", RED, "3 weeks"),
    ("P0: Admin Dashboard", "No content moderation, user management, or analytics visibility.\nRequired for IT Act compliance and operational control.", RED, "4 weeks"),
    ("P1: Payment Integration", "No revenue path. Razorpay for India + Stripe for international.\nScout/Team subscriptions are the primary revenue stream.", ORANGE, "3 weeks"),
    ("P1: Performance Sprint", "17 useState hooks in Discover, 25+ unoptimized FlatLists.\n666-line monolith components. Will fail on low-end Android.", ORANGE, "2 weeks"),
]

for i, (title, desc, color, effort) in enumerate(priorities):
    y = Inches(1.5 + i * 1.1)
    card = add_shape(slide, Inches(0.6), y, Inches(10.5), Inches(0.95), BG_CARD, color)
    add_text(slide, Inches(0.8), y + Inches(0.08), Inches(3), Inches(0.4), title, font_size=14, color=color, bold=True)
    add_text(slide, Inches(0.8), y + Inches(0.45), Inches(8), Inches(0.45), desc, font_size=10, color=LIGHT_GRAY)
    # Effort badge
    badge = add_shape(slide, Inches(11.3), y + Inches(0.2), Inches(1.5), Inches(0.5), color)
    add_text(slide, Inches(11.3), y + Inches(0.22), Inches(1.5), Inches(0.45), effort, font_size=12, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — THE ASK
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide("The Ask", "Seed round to scale India's sports talent platform")

add_text(slide, Inches(0.6), Inches(1.8), Inches(5), Inches(0.7), "Raising: Rs 2 Crore ($240K)", font_size=36, color=GREEN, bold=True)
add_text(slide, Inches(0.6), Inches(2.6), Inches(5), Inches(0.5), "Pre-money valuation: Rs 10 Crore ($1.2M)", font_size=18, color=LIGHT_GRAY)

use_data = [
    ["Category", "Allocation", "% of Raise", "Purpose"],
    ["Engineering", "Rs 60L ($72K)", "30%", "2 full-stack devs + 1 ML engineer for 12 months"],
    ["Growth & Marketing", "Rs 50L ($60K)", "25%", "User acquisition: college campaigns, sports events, digital ads"],
    ["Infrastructure", "Rs 20L ($24K)", "10%", "Supabase Pro, Claude API, CDN, monitoring"],
    ["Operations", "Rs 30L ($36K)", "15%", "Legal, compliance, office, travel for partnerships"],
    ["Content & Community", "Rs 20L ($24K)", "10%", "Athlete spotlight program, sports content creation"],
    ["Reserve", "Rs 20L ($24K)", "10%", "Buffer for unexpected costs and opportunities"],
]
add_table(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(3), use_data)

add_shape(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.5), BG_CARD, GREEN)
add_text(slide, Inches(0.8), Inches(6.62), Inches(11), Inches(0.45),
         "Target: 50K athletes + 100 paying scouts/teams within 12 months of funding  |  Path to profitability: Month 18",
         font_size=14, color=GREEN, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — CONTACT
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Large green accent
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(2), Inches(7), Inches(7))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x15, 0x30, 0x15)
circle.line.fill.background()

add_text(slide, Inches(1), Inches(1.5), Inches(8), Inches(1), "OnlyKrida", font_size=56, color=GREEN, bold=True)
add_text(slide, Inches(1), Inches(2.5), Inches(8), Inches(0.6), "Making Every Athlete Discoverable", font_size=24, color=WHITE)
add_text(slide, Inches(1), Inches(3.5), Inches(8), Inches(0.5), "India's First AI-Powered Sports Talent Discovery Platform", font_size=18, color=LIGHT_GRAY)

# Contact info
contact_items = [
    "Website: onlykrida.com",
    "Email: tanirudh127@gmail.com",
    "Platform: React Native + Supabase + Claude AI",
    "Status: MVP Complete — Ready to Scale",
]
add_bullet_list(slide, Inches(1), Inches(4.5), Inches(6), Inches(2), contact_items, font_size=16, color=LIGHT_GRAY, bullet_color=GREEN)

add_text(slide, Inches(1), Inches(6.5), Inches(6), Inches(0.4), "April 2026  |  Confidential", font_size=12, color=MUTED)


# ── Save ──────────────────────────────────────────────────────────────────────
output_path = "/Users/anirudhtumuluru/onlysports-platform/OnlyKrida_Complete_Deck.pptx"
prs.save(output_path)
print(f"\n{'='*60}")
print(f"  OnlyKrida Complete Deck generated successfully!")
print(f"  Saved to: {output_path}")
print(f"  Slides: {len(prs.slides)}")
print(f"{'='*60}\n")
